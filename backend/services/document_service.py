import os
import pymupdf as fitz
import docx
from pptx import Presentation
import pandas as pd
from typing import List, Dict, Any

class DocumentService:
    @classmethod
    def extract_text_and_pages(cls, file_path: str) -> List[Dict[str, Any]]:
        """
        Universal multi-format parser for educational materials:
        PDF, Word (DOCX), PowerPoint (PPTX), Text (TXT, MD), CSV, Excel (XLSX).
        Returns list of pages: [{"page_number": int, "text": str}]
        """
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            return cls._extract_pdf(file_path)
        elif ext in [".docx", ".doc"]:
            return cls._extract_docx(file_path)
        elif ext in [".pptx", ".ppt"]:
            return cls._extract_pptx(file_path)
        elif ext in [".txt", ".md", ".rtf"]:
            return cls._extract_txt(file_path)
        elif ext in [".xlsx", ".xls", ".csv"]:
            return cls._extract_tabular(file_path, ext)
        else:
            # Fallback to plain text
            return cls._extract_txt(file_path)

    @classmethod
    def _extract_pdf(cls, file_path: str) -> List[Dict[str, Any]]:
        pages_data = []
        doc = fitz.open(file_path)
        for page_idx in range(len(doc)):
            page = doc[page_idx]
            text = page.get_text("text").strip()
            if text:
                pages_data.append({
                    "page_number": page_idx + 1,
                    "text": text
                })
        doc.close()
        return pages_data or [{"page_number": 1, "text": "مستند PDF فارغ أو ممسوح ضوئياً."}]

    @classmethod
    def _extract_docx(cls, file_path: str) -> List[Dict[str, Any]]:
        doc = docx.Document(file_path)
        pages_data = []
        current_page_text = []
        page_num = 1
        word_count = 0

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            current_page_text.append(text)
            word_count += len(text.split())

            # Split into virtual pages every ~300 words or on page breaks
            if word_count >= 300:
                pages_data.append({
                    "page_number": page_num,
                    "text": "\n".join(current_page_text)
                })
                page_num += 1
                current_page_text = []
                word_count = 0

        # Extract tables
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_text:
                    table_rows.append(row_text)
            if table_rows:
                current_page_text.append("\n".join(table_rows))

        if current_page_text:
            pages_data.append({
                "page_number": page_num,
                "text": "\n".join(current_page_text)
            })

        return pages_data or [{"page_number": 1, "text": "مستند Word فارغ."}]

    @classmethod
    def _extract_pptx(cls, file_path: str) -> List[Dict[str, Any]]:
        prs = Presentation(file_path)
        pages_data = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_texts = []
            # Extract slide shapes and text frames
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
            
            # Extract slide notes if any
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[ملاحظات الشريحة: {notes}]")

            if slide_texts:
                pages_data.append({
                    "page_number": slide_idx + 1,
                    "text": f"--- [شريحة عرض {slide_idx + 1}] ---\n" + "\n".join(slide_texts)
                })

        return pages_data or [{"page_number": 1, "text": "عرض PowerPoint فارغ."}]

    @classmethod
    def _extract_txt(cls, file_path: str) -> List[Dict[str, Any]]:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
        except UnicodeDecodeError:
            with open(file_path, "r", encoding="latin-1") as f:
                content = f.read()

        paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
        pages_data = []
        current_chunk = []
        word_count = 0
        page_num = 1

        for para in paragraphs:
            current_chunk.append(para)
            word_count += len(para.split())
            if word_count >= 350:
                pages_data.append({
                    "page_number": page_num,
                    "text": "\n\n".join(current_chunk)
                })
                page_num += 1
                current_chunk = []
                word_count = 0

        if current_chunk:
            pages_data.append({
                "page_number": page_num,
                "text": "\n\n".join(current_chunk)
            })

        return pages_data or [{"page_number": 1, "text": content or "ملف نصي فارغ."}]

    @classmethod
    def _extract_tabular(cls, file_path: str, ext: str) -> List[Dict[str, Any]]:
        if ext == ".csv":
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)

        csv_text = df.to_string(index=False)
        return [{
            "page_number": 1,
            "text": f"--- [جدول بيانات تعليمي: {len(df)} سجل] ---\n" + csv_text[:10000]
        }]

    @classmethod
    def chunk_document(cls, pages_data: List[Dict[str, Any]], chunk_size: int = 250, overlap: int = 50) -> List[Dict[str, Any]]:
        chunks = []
        chunk_id = 1
        
        for page in pages_data:
            page_num = page["page_number"]
            text = page["text"]
            words = text.split()
            
            if len(words) <= chunk_size:
                chunks.append({
                    "chunk_id": f"c_{chunk_id}",
                    "page_number": page_num,
                    "text": text,
                    "word_count": len(words)
                })
                chunk_id += 1
            else:
                for i in range(0, len(words), chunk_size - overlap):
                    chunk_words = words[i:i + chunk_size]
                    if not chunk_words:
                        continue
                    chunk_text = " ".join(chunk_words)
                    chunks.append({
                        "chunk_id": f"c_{chunk_id}",
                        "page_number": page_num,
                        "text": chunk_text,
                        "word_count": len(chunk_words)
                    })
                    chunk_id += 1
        return chunks
